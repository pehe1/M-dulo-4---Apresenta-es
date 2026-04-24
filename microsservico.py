from fastapi import FastAPI, Form, HTTPException, File, UploadFile
from fastapi.responses import HTMLResponse, FileResponse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from google import genai
from google.genai import types
import requests
import os
import re
from typing import Optional
from dotenv import load_dotenv

# Carrega as variáveis de ambiente escondidas no arquivo .env
load_dotenv()

app = FastAPI()


# ==========================================
# FUNÇÃO DO MEMBRO 1: INGESTÃO DO GITHUB
# ==========================================
def extrair_commits_do_github(url: str) -> str:
    """
    Recebe uma URL do GitHub, acessa a API pública e retorna os últimos 10 commits em formato de texto.
    """
    # Usa Expressão Regular (Regex) para extrair o 'owner' e o 'repo' da URL
    match = re.search(r"github\.com/([^/]+)/([^/]+)", url)
    if not match:
        raise ValueError("URL do GitHub em formato inválido. Use: https://github.com/usuario/repositorio")

    owner, repo = match.groups()

    # Monta a URL da API do GitHub
    api_url = f"https://api.github.com/repos/{owner}/{repo}/commits"

    # Limpa o .git se o usuário colar o link de clone
    if repo.endswith('.git'):
        repo = repo[:-4]  # Remove as últimas 4 letras (".git")

    # Monta a URL da API do GitHub (repetido)
    api_url = f"https://api.github.com/repos/{owner}/{repo}/commits"

    # Faz a requisição na internet
    resposta = requests.get(api_url)

    # Valida se o repositório existe e é público
    if resposta.status_code == 404:
        raise ValueError("Repositório não encontrado. Verifique se a URL está correta e se o repositório é público.")
    elif resposta.status_code != 200:
        raise ValueError(f"Erro ao acessar o GitHub: {resposta.status_code}")

    # Extrai as mensagens dos últimos 10 commits
    dados_commits = resposta.json()
    historico_texto = "Histórico recente de atualizações do código:\n"

    for item in dados_commits[:10]:  # Limita aos 10 mais recentes
        mensagem = item.get("commit", {}).get("message", "Sem mensagem")
        autor = item.get("commit", {}).get("author", {}).get("name", "Desconhecido")
        # Remove quebras de linha das mensagens de commit para não confundir a IA
        mensagem_limpa = mensagem.replace('\n', ' | ')
        historico_texto += f"- {mensagem_limpa} (Autor: {autor})\n"

    return historico_texto


# ==========================================
# 1. FRONT-END: INGESTÃO UNIFICADA
# ==========================================
@app.get("/", response_class=HTMLResponse)
async def get_form():
    return """
    <html>
        <head>
            <title>Gerador de Apresentações com IA</title>
            <style>
                body { font-family: Arial, sans-serif; padding: 40px; background-color: #f4f4f9; }
                .container { background: white; padding: 30px; border-radius: 8px; box-shadow: 0 4px 8px rgba(0,0,0,0.1); max-width: 600px; margin: auto; }
                input[type=url], textarea, input[type=file] { width: 100%; padding: 10px; margin-top: 5px; margin-bottom: 20px; border: 1px solid #ccc; border-radius: 4px; box-sizing: border-box; }
                button { padding: 12px 20px; background: #6200ea; color: white; border: none; border-radius: 4px; cursor: pointer; font-size: 16px; font-weight: bold; width: 100%; }
                button:hover { background: #3700b3; }
                .tag-ia { display: inline-block; background: #e0f7fa; color: #00838f; padding: 4px 8px; border-radius: 12px; font-size: 12px; margin-bottom: 15px; font-weight: bold; }
                hr { border: 0; border-top: 1px solid #eee; margin: 20px 0; }
                label { font-weight: bold; color: #333; }
            </style>
        </head>
        <body>
            <div class="container">
                <h2>Gerador de Sprint Review 🚀</h2>
                <span class="tag-ia">✨ Powered by Gemini AI</span>
                <p>Escolha a origem dos dados para gerar a sua apresentação:</p>

                <form action="/gerar-pptx" method="post" enctype="multipart/form-data" onsubmit="document.getElementById('btn-gerar').innerText = 'Processando com IA... ⏳'; document.getElementById('btn-gerar').style.background = '#888';">

                    <label>Nova US: Análise de Diagrama Arquitetura</label>
                    <input type="file" name="imagem_diagrama" accept="image/*">
                    <hr>

                    <label>Nova US: Formatar Solução de Backend</label>
                    <textarea name="codigo_backend" rows="4" placeholder="Cole o código do backend aqui..."></textarea>
                    <hr>

                    <label>Opção 1: Extrair direto do Código</label>
                    <input type="url" name="url_github" placeholder="Ex: https://github.com/usuario/repositorio">
                    <hr>

                    <label>Opção 2: Colar anotações manualmente</label>
                    <textarea name="texto_bruto" rows="5" placeholder="Cole aqui o resumo da reunião, notas da Sprint, etc..."></textarea>
                    <hr>

                    <label for="tom_voz" style="display:block; margin-bottom: 5px; font-weight: bold;">Escolha o tom de voz:</label>
                    <select name="tom_voz" id="tom_voz" style="width: 100%; padding: 10px; margin-bottom: 20px; border: 1px solid #ccc; border-radius: 4px;">
                        <option value="formal">👔 Executivo/Formal (Foco em resultados e seriedade)</option>
                        <option value="persuasivo">🚀 Persuasivo/Entusiasta (Foco em engajamento e conquistas)</option>
                    </select>

                    <button type="submit" id="btn-gerar">Processar com IA e Gerar PPTX</button>
                </form>
            </div>
        </body>
    </html>
    """


# ==========================================
# 2. BACK-END: ROTEAMENTO INTELIGENTE
# ==========================================
# Note que agora usamos Optional[str] e Form(None) para não obrigar o usuário a preencher os dois
@app.post("/gerar-pptx")
async def gerar_pptx(
        url_github: Optional[str] = Form(None),
        texto_bruto: Optional[str] = Form(None),
        tom_voz: str = Form("formal"),
        codigo_backend: Optional[str] = Form(None),  # NOVO PARA US 1
        imagem_diagrama: UploadFile = File(None)  # NOVO PARA US 2
):
    chave = os.environ.get("GEMINI_API_KEY")
    cliente = genai.Client(api_key=chave)
    prs = Presentation()
    file_path = "Sprint_Review_Inteligente.pptx"

    # LÓGICA DE DECISÃO: De onde vem a informação?

    # ===============================================
    # IMPLEMENTAÇÃO: US 2 (DIAGRAMA DE ARQUITETURA)
    # ===============================================
    if imagem_diagrama and imagem_diagrama.filename:
        imagem_bytes = await imagem_diagrama.read()
        prompt_diagrama = "Você é um Arquiteto de Software. Analise este diagrama de arquitetura. Descreva o fluxo de dados principal em 3 a 5 tópicos curtos. Regra estrita: Sem asteriscos no início da frase e vá direto ao ponto."

        resposta_ia = cliente.models.generate_content(
            model='gemini-2.5-flash',
            contents=[
                types.Part.from_bytes(data=imagem_bytes, mime_type=imagem_diagrama.content_type),
                prompt_diagrama
            ]
        )

        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Análise da Arquitetura"
        corpo_texto = slide.placeholders[1].text_frame
        corpo_texto.paragraphs[0].text = "Fluxo de Dados Identificado:"

        for topico in resposta_ia.text.strip().split('\n'):
            if topico.strip():
                p = corpo_texto.add_paragraph()
                p.text = topico.strip()
                p.level = 1

    # ===============================================
    # IMPLEMENTAÇÃO: US 1 (CÓDIGO BACKEND)
    # ===============================================
    elif codigo_backend and codigo_backend.strip():
        prompt_codigo = f"Explique o objetivo deste código de backend em 1 única frase curta:\n\n{codigo_backend}"
        resposta_ia = cliente.models.generate_content(model='gemini-2.5-flash', contents=prompt_codigo)

        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Layout 5 é apenas Título
        slide.shapes.title.text = "Solução Backend"

        # Adiciona a explicação da IA no topo
        caixa_exp = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
        caixa_exp.text_frame.text = resposta_ia.text.strip()

        # Desenha o "Bloco de Código" (Retângulo preto com fonte Monoespaçada branca)
        caixa_cod = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(4))
        fill = caixa_cod.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(30, 30, 30)

        tf = caixa_cod.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = codigo_backend
        p.font.name = 'Courier New'
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(255, 255, 255)

    # ===============================================
    # IMPLEMENTAÇÃO ANTIGA: GITHUB OU MANUAL
    # ===============================================
    elif url_github or (texto_bruto and texto_bruto.strip()):
        if url_github:
            try:
                dados_entrada = extrair_commits_do_github(url_github)
                subtitulo_slide = f"Análise automatizada do repositório\n{url_github}"
            except ValueError as e:
                raise HTTPException(status_code=400, detail=str(e))

        elif texto_bruto and texto_bruto.strip():
            dados_entrada = f"Anotações da equipe:\n{texto_bruto}"
            subtitulo_slide = "Resumo gerado a partir de anotações manuais"

        # INGESTÃO CONCLUÍDA! Agora passa a bola para o Gemini (Membro 2)

        # Define instruções de tom baseadas na escolha do usuário
        instrucoes_tom = {
            "formal": "Mantenha um tom profissional, executivo e focado em resultados e métricas.",
            "persuasivo": "Mantenha um tom entusiasta, inspirador e focado em engajamento e conquistas da equipe."
        }

        prompt = f"""
        Você é um Tech Lead analisando informações de uma equipe de desenvolvimento.
        Traduza as seguintes informações em 3 a 5 tópicos profissionais e curtos para serem apresentados em um slide de Sprint Review para stakeholders.

        {instrucoes_tom.get(tom_voz, instrucoes_tom['formal'])}

        Regras estritas: 
        - Retorne APENAS os tópicos, um por linha.
        - NÃO use asteriscos, números, hífens ou marcadores no início da frase.
        - Vá direto ao ponto.

        Informações:
        {dados_entrada}
        """

        resposta_ia = cliente.models.generate_content(
            model='gemini-2.5-flash',
            contents=prompt,
        )

        topicos_processados = resposta_ia.text.strip().split('\n')

        # GERAÇÃO DO ARQUIVO PPTX (Membro 3)
        slide_capa = prs.slides.add_slide(prs.slide_layouts[0])
        slide_capa.shapes.title.text = "Sprint Review"
        slide_capa.placeholders[1].text = subtitulo_slide

        slide_conteudo = prs.slides.add_slide(prs.slide_layouts[1])
        slide_conteudo.shapes.title.text = "Principais Entregas"
        corpo_texto = slide_conteudo.placeholders[1].text_frame

        p0 = corpo_texto.paragraphs[0]
        p0.text = "Destaques da iteração:"

        for topico in topicos_processados:
            texto_limpo = topico.strip()
            if texto_limpo:
                p = corpo_texto.add_paragraph()
                p.text = texto_limpo
                p.level = 1
    else:
        raise HTTPException(status_code=400, detail="Você precisa preencher alguma das opções.")

    prs.save(file_path)

    return FileResponse(
        path=file_path,
        media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation',
        filename=file_path
    )